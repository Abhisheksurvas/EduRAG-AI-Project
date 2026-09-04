"""
rag_service.py — Production-ready RAG engine for EduRAG AI

Provides:
  - Text extraction from PDF, DOCX, PPTX, TXT, MD, CSV
  - Document chunking with overlap
  - Vector embedding generation via sentence-transformers
  - Cosine-similarity semantic retrieval (replaces keyword BM25)
  - Strict source attribution and hallucination prevention
"""

from __future__ import annotations

import base64
import io
import logging
import os
import re
import uuid
from datetime import datetime, timezone
from functools import lru_cache
from typing import Optional

# Windows without Developer Mode cannot create the symlinks HuggingFace uses
# for its cache; disable them so the embedding model downloads/caches via
# plain file copies instead of hanging on a symlink operation.
os.environ.setdefault("HF_HUB_DISABLE_SYMLINKS", "1")

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
CHUNK_WORDS: int = 250
CHUNK_OVERLAP: int = 50
SIMILARITY_THRESHOLD: float = 0.35
DEFAULT_TOP_K: int = 5
EMBEDDING_MODEL: str = "all-MiniLM-L6-v2"

TOKEN_PATTERN = re.compile(r"[a-zA-Z0-9]{2,}")


# ---------------------------------------------------------------------------
# Embedding subsystem  (lazy-loaded so the server starts even without GPU)
# ---------------------------------------------------------------------------

_embedding_model = None          # SentenceTransformer instance, loaded on first use
_embedding_available: Optional[bool] = None  # None = not yet tried


def _load_embedding_model():
    """Load the sentence-transformer model once and cache it."""
    global _embedding_model, _embedding_available

    if _embedding_available is not None:
        return _embedding_available

    try:
        from sentence_transformers import SentenceTransformer  # type: ignore
        _embedding_model = SentenceTransformer(EMBEDDING_MODEL)
        _embedding_available = True
        logger.info("[RAG] Sentence-transformer model loaded: %s", EMBEDDING_MODEL)
    except Exception as exc:  # pragma: no cover
        logger.warning("[RAG] sentence-transformers unavailable, falling back to keyword search: %s", exc)
        _embedding_available = False

    return _embedding_available


def generate_embeddings(text: str) -> list[float]:
    """
    Return a normalised float32 embedding vector for *text*.

    Falls back to an empty list when sentence-transformers is not installed,
    allowing the system to degrade gracefully to keyword search.
    """
    if not text or not text.strip():
        return []

    if not _load_embedding_model():
        return []

    try:
        import numpy as np  # type: ignore
        vector = _embedding_model.encode(text.strip(), normalize_embeddings=True)
        return vector.tolist()
    except Exception as exc:  # pragma: no cover
        logger.error("[RAG] Embedding generation failed: %s", exc)
        return []


def generate_embeddings_batch(texts: list[str]) -> list[list[float]]:
    """
    Return normalised float32 embedding vectors for a batch of *texts* in a
    single vectorised call.  Encoding a list at once is dramatically faster
    than calling :func:`generate_embeddings` per chunk (GPU/CPU batching and
    amortised overhead), which is what makes document indexing quick.
    """
    if not texts:
        return []

    cleaned: list[str] = [t.strip() for t in texts if t and t.strip()]
    if not cleaned:
        return [[] for _ in texts]

    if not _load_embedding_model():
        return [[] for _ in texts]

    try:
        vectors = _embedding_model.encode(
            cleaned,
            normalize_embeddings=True,
            batch_size=64,
            show_progress_bar=False,
        )
        return [v.tolist() for v in vectors]
    except Exception as exc:  # pragma: no cover
        logger.error("[RAG] Batch embedding generation failed: %s", exc)
        return [[] for _ in texts]


def warm_up_embeddings() -> None:
    """
    Load the embedding model up front (e.g. at server startup) so the first
    document upload doesn't pay the lazy-load cost.  A tiny dummy encode
    triggers model download/instantiation without embedding a real document.
    Failures are non-fatal — the system degrades to keyword search.
    """
    try:
        if _load_embedding_model():
            _embedding_model.encode(["warmup"], normalize_embeddings=True)
            logger.info("[RAG] Embedding model warmed up and ready for fast indexing.")
    except Exception as exc:  # pragma: no cover
        logger.warning("[RAG] Embedding warm-up skipped: %s", exc)


# ---------------------------------------------------------------------------
# Text-extraction helpers
# ---------------------------------------------------------------------------

def tokens(value: str) -> list[str]:
    return TOKEN_PATTERN.findall(value.lower())


def extract_text(filename: str, encoded_content: str) -> list[dict]:
    """Decode base64 content and extract pages."""
    raw = base64.b64decode(encoded_content)
    return extract_text_from_bytes(filename, raw)


IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "bmp", "tiff", "tif", "gif", "webp"}


def _run_ocr_on_image(image) -> str:
    """
    Run OCR on a single PIL image and return the recognized text.

    Tries EasyOCR first (bundles its own engine, reuses the existing torch
    dependency, no external binary required), then falls back to pytesseract
    if Tesseract-OCR is installed on the host. Returns "" if no OCR engine is
    available or recognition fails.
    """
    # EasyOCR — self-contained, works without system Tesseract.
    try:
        import easyocr  # type: ignore

        # Lazily build a reader for English; reuse across calls via lru_cache.
        reader = _get_easyocr_reader()
        if reader is not None:
            results = reader.readtext(_pil_to_array(image), detail=0, paragraph=True)
            text = "\n".join(str(line).strip() for line in results if str(line).strip())
            if text.strip():
                return text
    except Exception as exc:  # pragma: no cover - optional engine
        logger.warning("[RAG] EasyOCR failed: %s", exc)

    # pytesseract fallback — requires Tesseract-OCR binary on the host.
    try:
        import pytesseract  # type: ignore

        text = pytesseract.image_to_string(image).strip()
        if text:
            return text
    except Exception as exc:  # pragma: no cover - optional engine
        logger.warning("[RAG] pytesseract failed: %s", exc)

    return ""


@lru_cache(maxsize=1)
def _get_easyocr_reader():
    try:
        import easyocr  # type: ignore

        return easyocr.Reader(["en"], gpu=False)
    except Exception as exc:  # pragma: no cover
        logger.warning("[RAG] Could not initialise EasyOCR reader: %s", exc)
        return None


def _pil_to_array(image):
    try:
        import numpy as np  # type: ignore

        return np.array(image.convert("RGB"))
    except Exception:  # pragma: no cover
        return image


def _ocr_pdf_pages(raw: bytes) -> list[dict]:
    """
    Render each PDF page to an image and OCR it. Used as a fallback when a PDF
    has no extractable (text-layer) content — e.g. a scanned document.
    """
    try:
        import pypdfium2  # type: ignore

        pdf = pypdfium2.PdfDocument(io.BytesIO(raw))
        pages: list[dict] = []
        for idx in range(len(pdf)):
            try:
                page = pdf[idx]
                bitmap = page.render(scale=2.0)
                pil_image = bitmap.to_pil()
                text = _run_ocr_on_image(pil_image).strip()
                pages.append({"page": idx + 1, "text": text})
            except Exception as exc:  # pragma: no cover
                logger.warning("[RAG] OCR failed for PDF page %s: %s", idx + 1, exc)
                pages.append({"page": idx + 1, "text": ""})
        return pages
    except Exception as exc:  # pragma: no cover
        logger.warning("[RAG] Could not render PDF for OCR: %s", exc)
        return []


def extract_text_from_bytes(filename: str, raw: bytes) -> list[dict]:
    """Return extracted pages from raw bytes — no base64 involved."""
    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

    if extension in {"txt", "md", "csv"}:
        return [{"page": 1, "text": raw.decode("utf-8", errors="replace")}]

    if extension in IMAGE_EXTENSIONS:
        return _extract_text_from_image(raw)

    if extension == "pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError as exc:
            raise ValueError("PDF processing requires pypdf on the server.") from exc
        reader = PdfReader(io.BytesIO(raw))
        pages = [
            {"page": idx + 1, "text": page.extract_text() or ""}
            for idx, page in enumerate(reader.pages)
        ]

        # Scanned/image PDFs often have an empty text layer. Fall back to OCR
        # so the document can still be indexed instead of being rejected.
        if sum(len(p.get("text", "").strip()) for p in pages) < 40:
            logger.info("[RAG] PDF has no text layer — attempting OCR fallback")
            ocr_pages = _ocr_pdf_pages(raw)
            if sum(len(p.get("text", "").strip()) for p in ocr_pages) > 0:
                pages = ocr_pages
        return pages

    if extension in {"doc", "docx"}:
        if extension == "doc":
            raise ValueError(
                "Legacy .doc files are not supported. Please save as .docx and re-upload."
            )
        try:
            from docx import Document  # type: ignore
        except ImportError as exc:
            raise ValueError("DOCX processing requires python-docx on the server.") from exc
        document = Document(io.BytesIO(raw))
        return [{"page": 1, "text": "\n".join(p.text for p in document.paragraphs)}]

    if extension in {"ppt", "pptx"}:
        if extension == "ppt":
            raise ValueError(
                "Legacy .ppt files are not supported. Please save as .pptx and re-upload."
            )
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise ValueError("PPTX processing requires python-pptx on the server.") from exc
        presentation = Presentation(io.BytesIO(raw))
        pages = []
        for idx, slide in enumerate(presentation.slides):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            pages.append({"page": idx + 1, "text": text})
        return pages

    raise ValueError(
        f"Unsupported file type '.{extension}'. Upload PDF, DOCX, PPTX, TXT, MD, or CSV."
    )


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_pages(pages: list[dict], material: dict) -> list[dict]:
    """
    Split pages into overlapping word-window chunks and attach an embedding
    vector to each chunk.  Embedding is stored under the key ``'embedding'``.

    Chunk texts are gathered first and embedded in a single batched call via
    :func:`generate_embeddings_batch` for speed.
    """
    chunks: list[dict] = []
    parts: list[str] = []
    for page_data in pages:
        words = str(page_data.get("text", "")).split()
        for start in range(0, len(words), CHUNK_WORDS - CHUNK_OVERLAP):
            part = " ".join(words[start : start + CHUNK_WORDS]).strip()
            if len(tokens(part)) < 8:
                continue

            parts.append(part)
            chunks.append(
                {
                    "id": f"chunk-{uuid.uuid4().hex}",
                    "materialId": material["id"],
                    "studentId": material["studentId"],
                    "documentName": material["name"],
                    "page": page_data["page"],
                    "text": part,
                    "embedding": [],   # filled in by the batched call below
                }
            )

    if parts:
        # Generate all embeddings in one vectorised batch — far quicker than
        # encoding each chunk individually.
        embeddings = generate_embeddings_batch(parts)
        for chunk, embedding in zip(chunks, embeddings):
            chunk["embedding"] = embedding

    return chunks


# ---------------------------------------------------------------------------
# Material creation
# ---------------------------------------------------------------------------

def create_material(body: dict) -> tuple[dict, list[dict]]:
    """
    Parse a base64-encoded upload request and return
    ``(material_record, chunks_with_embeddings)``.
    """
    filename = str(body.get("name", "")).strip()
    student_id = str(body.get("studentId", "")).strip()
    encoded_content = str(body.get("contentBase64", ""))
    if not filename or not student_id or not encoded_content:
        raise ValueError("studentId, name, and file content are required.")

    pages = extract_text(filename, encoded_content)
    _validate_extracted_pages(pages)

    material = _build_material_record(filename, student_id, body, pages)
    chunks = chunk_pages(pages, material)

    if not chunks:
        raise ValueError(
            "The uploaded file did not contain enough readable study material to index."
        )
    return material, chunks


def create_material_from_bytes(
    filename: str,
    raw: bytes,
    student_id: str,
    course: str = "Personal study material",
) -> tuple[dict, list[dict]]:
    """
    Create a material record and chunks from raw file bytes (no base64).
    Embeddings are generated and attached to every chunk.
    """
    filename = filename.strip()
    student_id = student_id.strip()
    if not filename or not student_id or not raw:
        raise ValueError("studentId, filename, and file content are required.")

    pages = extract_text_from_bytes(filename, raw)
    _validate_extracted_pages(pages)

    file_base64 = base64.b64encode(raw).decode("utf-8")
    material = _build_material_record(
        filename, student_id, {"course": course}, pages, file_base64=file_base64
    )
    chunks = chunk_pages(pages, material)

    if not chunks:
        raise ValueError(
            "The uploaded file did not contain enough readable study material to index."
        )
    return material, chunks


def _extract_text_from_image(raw: bytes) -> list[dict]:
    """OCR a standalone image file (scanned page, photo of notes, etc.)."""
    try:
        from PIL import Image  # type: ignore
    except ImportError as exc:
        raise ValueError(
            "Image OCR requires Pillow on the server. Install with 'pip install Pillow'."
        ) from exc
    try:
        image = Image.open(io.BytesIO(raw))
        text = _run_ocr_on_image(image).strip()
        return [{"page": 1, "text": text}]
    except Exception as exc:
        raise ValueError(f"Could not read the image file: {exc}")


def _validate_extracted_pages(pages: list[dict]) -> None:
    extracted_chars = sum(len(str(p.get("text", "")).strip()) for p in pages)
    # A document can be perfectly valid while containing only a title, a short
    # formula, or a small excerpt.  The old 40-character cut-off made those
    # PDFs fail at upload time even though their text was extracted correctly.
    # Only reject files for which neither the PDF text layer nor OCR yielded
    # any searchable content at all.
    if extracted_chars == 0:
        raise ValueError(
            "No text could be extracted (scanned/image PDF with no OCR support?). "
            "Install an OCR engine (pip install easyocr) and re-upload, or use a "
            "text-based document."
        )


def _build_material_record(
    filename: str, student_id: str, body: dict, pages: list[dict], file_base64: str | None = None
) -> dict:
    record = {
        "id": f"material-{uuid.uuid4().hex}",
        "studentId": student_id,
        "name": filename,
        "course": str(body.get("course", "Personal study material")),
        "type": filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt",
        "uploadedBy": "Student",
        "uploadedAt": datetime.now(timezone.utc).isoformat(),
        "pages": len(pages),
    }
    if file_base64:
        record["fileBase64"] = file_base64
    return record


# ---------------------------------------------------------------------------
# Cosine-similarity helpers
# ---------------------------------------------------------------------------

def _cosine_similarity(vec_a: list[float], vec_b: list[float]) -> float:
    """
    Compute cosine similarity between two pre-normalised vectors.
    Returns 0.0 for empty or mismatched vectors.
    """
    if not vec_a or not vec_b or len(vec_a) != len(vec_b):
        return 0.0
    try:
        import numpy as np  # type: ignore
        a = np.array(vec_a, dtype="float32")
        b = np.array(vec_b, dtype="float32")
        # Vectors from sentence-transformers are already unit-normalised;
        # dot product equals cosine similarity.
        score = float(np.dot(a, b))
        return max(0.0, min(1.0, score))
    except Exception:
        # Pure-Python fallback (slower but dependency-free)
        dot = sum(x * y for x, y in zip(vec_a, vec_b))
        norm_a = sum(x * x for x in vec_a) ** 0.5
        norm_b = sum(x * x for x in vec_b) ** 0.5
        if norm_a == 0 or norm_b == 0:
            return 0.0
        return max(0.0, min(1.0, dot / (norm_a * norm_b)))


# ---------------------------------------------------------------------------
# Semantic retrieval  (replaces keyword-based retrieve())
# ---------------------------------------------------------------------------

def semantic_retrieve(
    query: str,
    chunks: list[dict],
    top_k: int = DEFAULT_TOP_K,
    selected_material_ids: list[str] | None = None,
    threshold: float = SIMILARITY_THRESHOLD,
) -> list[dict]:
    """
    Retrieve the *top_k* most semantically relevant chunks using cosine
    similarity between the query embedding and stored chunk embeddings.

    Falls back to keyword-overlap scoring when embeddings are unavailable.

    Args:
        query:                  The user's question.
        chunks:                 All indexed chunks (each may have 'embedding').
        top_k:                  Maximum number of chunks to return.
        selected_material_ids:  If non-empty, restrict search to these materials.
        threshold:              Minimum similarity score; lower chunks are dropped.

    Returns:
        List of chunk dicts augmented with a ``'score'`` float field, sorted
        descending by score.  Returns an empty list when no chunks meet the
        threshold — the caller must handle the no-context case.
    """
    query = query.strip()
    if not query or not chunks:
        return []

    selected = set(selected_material_ids) if selected_material_ids else set()

    # Filter by material selection
    candidate_chunks = [
        c for c in chunks
        if not selected or c.get("materialId") in selected
    ]

    if not candidate_chunks:
        return []

    # Attempt vector search only when both the query and the candidate chunks
    # actually carry usable embeddings. If the chunks were indexed while the
    # embedding model was unavailable they have empty embeddings; in that case
    # _vector_search would skip every chunk and return nothing, so we must fall
    # back to keyword search instead of silently returning an empty result.
    query_embedding = generate_embeddings(query)
    chunks_with_embeddings = [c for c in candidate_chunks if c.get("embedding")]
    if query_embedding and chunks_with_embeddings:
        return _vector_search(query_embedding, candidate_chunks, top_k, threshold)

    # Graceful fallback: keyword overlap (used when embeddings unavailable)
    logger.info("[RAG] Embeddings unavailable — using keyword fallback for retrieval.")
    return _keyword_search(query, candidate_chunks, top_k)


def _vector_search(
    query_embedding: list[float],
    chunks: list[dict],
    top_k: int,
    threshold: float,
) -> list[dict]:
    """Score chunks by cosine similarity and return top-k above threshold."""
    scored: list[tuple[float, dict]] = []

    for chunk in chunks:
        chunk_embedding = chunk.get("embedding") or []
        if not chunk_embedding:
            # Chunk has no embedding stored — skip it in vector mode
            continue
        score = _cosine_similarity(query_embedding, chunk_embedding)
        if score >= threshold:
            scored.append((score, chunk))

    scored.sort(key=lambda item: item[0], reverse=True)
    return [
        {**chunk, "score": round(score, 4)}
        for score, chunk in scored[:top_k]
    ]


def _keyword_search(
    query: str,
    chunks: list[dict],
    top_k: int,
) -> list[dict]:
    """BM25-style keyword overlap search (fallback when no embeddings)."""
    from collections import Counter

    question_terms = Counter(tokens(query))
    if not question_terms:
        return []

    scored: list[tuple[float, dict]] = []
    for chunk in chunks:
        chunk_terms = Counter(tokens(str(chunk.get("text", ""))))
        overlap = sum(
            question_terms[t] * min(chunk_terms[t], 3) for t in question_terms
        )
        phrase_bonus = 3 if query.lower() in str(chunk.get("text", "")).lower() else 0
        total = overlap + phrase_bonus
        if total:
            scored.append((float(total), chunk))

    scored.sort(key=lambda item: item[0], reverse=True)
    return [
        {**chunk, "score": round(score, 4)}
        for score, chunk in scored[:top_k]
    ]


# ---------------------------------------------------------------------------
# Backward-compatible alias  (app.py still calls retrieve())
# ---------------------------------------------------------------------------

def retrieve(
    question: str,
    chunks: list[dict],
    selected_material_ids: list[str],
    limit: int = DEFAULT_TOP_K,
) -> list[dict]:
    """
    Drop-in replacement for the old keyword-based retrieve().
    Delegates to semantic_retrieve() which prefers vector search and falls
    back to keyword overlap automatically.
    """
    return semantic_retrieve(
        query=question,
        chunks=chunks,
        top_k=limit,
        selected_material_ids=selected_material_ids,
    )


# ---------------------------------------------------------------------------
# Answer generation  (strict grounding — no hallucination)
# ---------------------------------------------------------------------------

NO_CONTEXT_RESPONSE = (
    "I couldn't find this information in your available study materials."
)

LOW_RELEVANCE_RESPONSE = (
    "I couldn't find this information in your available study materials."
)


def extractive_answer(question: str, retrieved: list[dict]) -> str:
    """
    Build a grounded answer strictly from the retrieved chunks.

    Rules:
    - Never fabricate information not present in retrieved chunks.
    - If retrieved list is empty, return NO_CONTEXT_RESPONSE.
    - Each returned sentence must originate from a retrieved chunk.
    - Include source attribution (document name + page number).
    """
    if not retrieved:
        return NO_CONTEXT_RESPONSE

    question_terms = set(tokens(question))
    question_keywords = set()
    for term in question_terms:
        if len(term) > 2:
            question_keywords.add(term)
    question_stop = {"the", "is", "at", "which", "who", "what", "where", "when", "why", "how", "does", "do", "did", "can", "could", "would", "should", "may", "might", "shall", "will", "has", "have", "had", "been", "was", "were", "are", "am", "a", "an", "of", "to", "in", "for", "on", "with", "as", "by", "from", "that", "this", "it", "its"}
    question_keywords -= question_stop
    if not question_keywords:
        question_keywords = question_terms

    candidates: list[tuple[int, str, str, int]] = []

    for chunk in retrieved:
        doc_name = chunk.get("documentName", "Unknown document")
        page_num = chunk.get("page", "?")
        chunk_text = str(chunk.get("text", ""))
        for sentence in re.split(r'(?<=[.!?])\s+', chunk_text):
            cleaned = sentence.strip()
            if len(cleaned) < 30:
                continue
            sentence_terms = set(tokens(cleaned))
            if not sentence_terms:
                continue
            overlap = len(question_keywords.intersection(sentence_terms))
            overlap_ratio = overlap / len(sentence_terms) if sentence_terms else 0
            score = overlap * 2 + int(overlap_ratio * 3)
            if score > 0:
                candidates.append((score, cleaned, doc_name, page_num))

    candidates.sort(key=lambda item: item[0], reverse=True)

    seen: set[str] = set()
    unique: list[tuple[str, str, int]] = []
    for _, sentence, doc_name, page_num in candidates:
        normalized = sentence.lower().strip()
        if normalized in seen:
            continue
        seen.add(normalized)
        unique.append((sentence, doc_name, page_num))
        if len(unique) == DEFAULT_TOP_K:
            break

    if not unique:
        return LOW_RELEVANCE_RESPONSE

    lines = [
        f"- {sentence}\n  *(Source: {doc_name}, page {page_num})*"
        for sentence, doc_name, page_num in unique
    ]
    return "Based on your study materials:\n\n" + "\n\n".join(lines)


def format_sources(retrieved: list[dict]) -> list[dict]:
    """
    Return a deduplicated list of source references for the API response.
    Each entry contains the document name, page number, and similarity score.
    """
    seen: set[tuple] = set()
    sources: list[dict] = []
    for chunk in retrieved:
        key = (chunk.get("documentName"), chunk.get("page"))
        if key not in seen:
            seen.add(key)
            sources.append(
                {
                    "doc": chunk.get("documentName", "Unknown"),
                    "page": chunk.get("page", "?"),
                    "score": round(float(chunk.get("score", 0.0)), 4),
                }
            )
    return sources
